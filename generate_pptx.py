# -*- coding: utf-8 -*-
"""
generate_pptx.py
依赖：python-pptx
保存 slides_content.json 并与本脚本放在同一目录，然后运行：
    pip install python-pptx
    python generate_pptx.py
输出文件名由 slides_content.json 中 output_filename 指定（默认：三国两晋南北朝_第6课_AI赋能教学.pptx）
"""

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

CONTENT_FILE = "slides_content.json"

def set_slide_title(slide, title):
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    # set title font size
    for paragraph in title_placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(28)
            run.font.bold = True

def add_bullets(slide, bullets):
    # find a content placeholder or add textbox
    left = Inches(0.6)
    top = Inches(1.6)
    width = Inches(8.2)
    height = Inches(4.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.margin_bottom = Pt(6)
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = '微软雅黑'
        p.font.color.rgb = RGBColor(0, 0, 0)

def add_notes(prs_slide, notes_text):
    notes_slide = prs_slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes_text

def add_image_placeholder(slide, caption=None):
    # Inserts a simple rectangle with caption text as a visual placeholder
    left = Inches(6.5)
    top = Inches(1.6)
    width = Inches(3)
    height = Inches(3)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.text = caption or "图片占位（请用AI生成图像并替换）"
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(12)
            run.font.italic = True

def build_presentation(content):
    prs = Presentation()
    # set slide width/height if needed
    for s in content['slides']:
        slide_layout = prs.slide_layouts[5]  # blank with title to place content
        slide = prs.slides.add_slide(slide_layout)
        set_slide_title(slide, s.get('title', ''))
        bullets = s.get('bullets', [])
        if bullets:
            add_bullets(slide, bullets)
        # add image placeholder if prompt exists
        if s.get('image_prompt'):
            add_image_placeholder(slide, caption="图片占位（提示：%s）" % s.get('image_prompt')[:80])
        # add notes
        notes = s.get('notes', '')
        if notes:
            add_notes(slide, notes)
    # save
    out = content.get('output_filename', 'output.pptx')
    prs.save(out)
    return out

def main():
    if not os.path.exists(CONTENT_FILE):
        print("找不到 %s，请把 slides_content.json 与脚本放在同一目录。" % CONTENT_FILE)
        return
    with open(CONTENT_FILE, 'r', encoding='utf-8') as f:
        content = json.load(f)
    out = build_presentation(content)
    print("已生成 PPTX：%s" % out)
    print("提示：若要把AI生成的图片嵌入PPT，请先用附件中提供的 prompts 在你选择的图像生成工具生成图片，"
          "然后在 PPT 中以同名替换或手动插入图片到对应幻灯片。")
    print("建议：若需我帮助生成英/中版本的图像 prompts 或直接生成单张示意图，请告知。")

if __name__ == "__main__":
    main()